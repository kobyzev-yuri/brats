"""
FastAPI сервис для управления базой знаний (KB)
Адаптировано из ~/sql4A/src/api/main.py
"""

import sys
import os
import logging
import time
from pathlib import Path
from typing import Dict, Any, List, Optional

# Загружаем конфигурацию: сначала корень проекта (как у скриптов и start_all_services), затем kb-service (переопределение)
from dotenv import load_dotenv
_root = Path(__file__).resolve().parents[2]  # корень репо brats
_kb_dir = Path(__file__).resolve().parents[1]  # kb-service
load_dotenv(dotenv_path=_root / "config.env")
load_dotenv(dotenv_path=_kb_dir / "config.env", override=True)

from fastapi import Body, FastAPI, File, Form, HTTPException, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse
from fastapi.staticfiles import StaticFiles

from models.requests import (
    KBChunkCreate, KBChunkUpdate, KBSearchRequest, KBImportRequest,
    Category, TargetAudience,
)
from models.responses import (
    KBChunkResponse, KBSearchResponse, KBAddResponse, KBUpdateResponse,
    KBDeleteResponse, KBImportResponse, KBStatsResponse, KBListResponse, KBChunkListItem,
    HealthCheckResponse,
)
from services.kb_service import KBService
from utils.db import get_db_connection, ensure_pgvector_extension
from api.rag_endpoints import router as rag_router
from api.mmkb_endpoints import router as mmkb_router
from api.dlp_endpoints import router as dlp_router

# Настройка логирования
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[logging.StreamHandler()]
)
logger = logging.getLogger(__name__)

# Создание FastAPI приложения
app = FastAPI(
    title="KB Service API",
    description="API для управления базой знаний (KB) на основе PostgreSQL + pgvector",
    version="1.0.0",
    docs_url="/docs",
    redoc_url="/redoc"
)

# Настройка CORS
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # В продакшене указать конкретные домены
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Инициализация сервиса KB
kb_service = KBService()

# --- DLP (первыми, чтобы точно были зарегистрированы) ---
@app.get("/api/dlp/health")
@app.get("/dlp/health")  # дубль без /api на случай перехвата префикса
def api_dlp_health():
    """Проверка доступности DLP API."""
    return {"status": "ok", "service": "dlp"}


@app.post("/api/dlp/sanitize-text")
def api_dlp_sanitize_text(body: Dict[str, Any] = Body(None)):
    """Обезличивание текста (телефоны, email и т.д.) для n8n/LLM."""
    if body is None:
        body = {}
    try:
        from services.dlp_service import get_dlp_service
        text = body.get("text") or ""
        if not isinstance(text, str):
            text = str(text) if text is not None else ""
        dlp = get_dlp_service()
        sanitized = dlp._mask_text(text)
        return {"sanitized_text": sanitized}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


def _log_routes():
    """При старте проверить, что DLP зарегистрирован."""
    for r in app.routes:
        if hasattr(r, "path") and "dlp" in r.path:
            logger.info("DLP route registered: %s %s", getattr(r, "methods", ""), r.path)


# Подключаем RAG endpoints
app.include_router(rag_router)

# Подключаем мультимодальную KB (mmkb)
app.include_router(mmkb_router)

# DLP router (дублирует маршруты выше при успешной загрузке)
try:
    app.include_router(dlp_router)
except Exception as e:
    logger.warning("DLP router not loaded: %s", e)


# Каталог для загруженных изображений (блоки оператора)
UPLOADS_DIR = Path(__file__).resolve().parents[1] / "uploads"
UPLOADS_DIR.mkdir(exist_ok=True)
app.mount("/static", StaticFiles(directory=str(UPLOADS_DIR)), name="mmkb_static")


@app.get("/", response_model=Dict[str, str])
async def root():
    """
    Корневой эндпоинт
    """
    return {
        "message": "KB Service API работает",
        "version": "1.0.0",
        "docs": "/docs"
    }


@app.get("/health", response_model=HealthCheckResponse)
async def health_check():
    """
    Проверка здоровья системы
    """
    try:
        # Проверяем подключение к БД
        conn = await get_db_connection()
        await ensure_pgvector_extension(conn)
        pgvector_ok = True
        await conn.close()
        
        return HealthCheckResponse(
            status="healthy",
            database="connected",
            pgvector_extension=True,
            version="1.0.0"
        )
    except Exception as e:
        logger.error(f"Ошибка проверки здоровья: {e}")
        return HealthCheckResponse(
            status="unhealthy",
            database="disconnected",
            pgvector_extension=False,
            version="1.0.0"
        )


@app.post("/api/kb/search", response_model=KBSearchResponse)
async def search_kb(request: KBSearchRequest):
    """
    Семантический поиск в базе знаний
    """
    try:
        start_time = time.time()
        
        results = await kb_service.search(
            query=request.query,
            limit=request.limit,
            category=request.category.value if request.category else None,
            target_audience=request.target_audience.value if request.target_audience else None,
            priority=request.priority.value if request.priority else None,
            settlement_id=request.settlement_id,
            min_similarity=request.min_similarity,
            use_case=request.use_case,
            stage=request.stage,
            source=request.source,
        )
        
        search_time_ms = (time.time() - start_time) * 1000
        
        return KBSearchResponse(
            query=request.query,
            total_found=len(results),
            results=[
                KBChunkResponse(
                    id=r["id"],
                    content=r["content"],
                    metadata=r["metadata"],
                    version=r["version"],
                    similarity=r.get("similarity"),
                    created_at=r["created_at"],
                    updated_at=r["updated_at"],
                    last_updated=r["last_updated"],
                    is_active=r["is_active"]
                )
                for r in results
            ],
            search_time_ms=round(search_time_ms, 2)
        )
    except Exception as e:
        logger.error(f"Ошибка поиска в KB: {e}")
        import traceback
        logger.error(traceback.format_exc())
        raise HTTPException(status_code=500, detail=f"Ошибка поиска: {str(e)}")


@app.post("/api/kb/add", response_model=KBAddResponse)
async def add_chunk(request: KBChunkCreate):
    """
    Добавление нового chunk в KB
    """
    try:
        chunk_id = await kb_service.add_chunk(
            content=request.content,
            category=request.category.value,
            target_audience=request.target_audience.value,
            priority=request.priority.value,
            subcategory=request.subcategory,
            tags=request.tags,
            source=request.source,
            version=request.version,
            related_links=request.related_links,
            use_case=request.use_case,
            stage=request.stage,
            settlement_id=request.settlement_id,
            metadata=request.metadata
        )
        
        return KBAddResponse(
            success=True,
            chunk_id=chunk_id,
            message=f"Chunk успешно добавлен (id={chunk_id})"
        )
    except Exception as e:
        logger.error(f"Ошибка добавления chunk: {e}")
        import traceback
        logger.error(traceback.format_exc())
        raise HTTPException(status_code=500, detail=f"Ошибка добавления: {str(e)}")


@app.get("/api/kb/stats", response_model=KBStatsResponse)
async def get_kb_stats():
    """
    Получение статистики по KB
    """
    try:
        stats = await kb_service.get_stats()
        
        # Преобразуем last_updated в datetime если это строка или None
        last_updated = stats.get("last_updated")
        if last_updated and isinstance(last_updated, str):
            from datetime import datetime
            try:
                last_updated = datetime.fromisoformat(last_updated.replace('Z', '+00:00'))
            except:
                last_updated = None
        
        return KBStatsResponse(
            total_chunks=stats.get("total_chunks", 0),
            active_chunks=stats.get("active_chunks", 0),
            chunks_by_category=stats.get("chunks_by_category", {}),
            chunks_by_target_audience=stats.get("chunks_by_target_audience", {}),
            chunks_by_priority=stats.get("chunks_by_priority", {}),
            last_updated=last_updated
        )
    except Exception as e:
        logger.error(f"Ошибка получения статистики: {e}")
        import traceback
        logger.error(traceback.format_exc())
        raise HTTPException(status_code=500, detail=f"Ошибка получения статистики: {str(e)}")


@app.get("/api/kb/list", response_model=KBListResponse)
async def list_chunks(limit: int = 2000, source: Optional[str] = None):
    """
    Список активных chunk (id + title) для выбора при редактировании.
    source: опционально — фильтр по источнику (подстрока).
    """
    try:
        chunks = await kb_service.list_chunks(limit=limit, source=source)
        return KBListResponse(
            chunks=[KBChunkListItem(id=c["id"], title=c["title"]) for c in chunks]
        )
    except Exception as e:
        logger.error(f"Ошибка list_chunks: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@app.get("/api/kb/sources")
async def get_sources():
    """Список уникальных источников (metadata.source) для фильтра."""
    try:
        return await kb_service.list_sources()
    except Exception as e:
        logger.error(f"Ошибка get_sources: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@app.get("/api/kb/{chunk_id}", response_model=KBChunkResponse)
async def get_chunk(chunk_id: int):
    """
    Получение chunk по ID
    """
    try:
        chunk = await kb_service.get_chunk(chunk_id)
        if not chunk:
            raise HTTPException(status_code=404, detail=f"Chunk {chunk_id} не найден")
        
        return KBChunkResponse(
            id=chunk["id"],
            content=chunk["content"],
            metadata=chunk["metadata"],
            version=chunk["version"],
            created_at=chunk["created_at"],
            updated_at=chunk["updated_at"],
            last_updated=chunk["last_updated"],
            is_active=chunk["is_active"]
        )
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Ошибка получения chunk {chunk_id}: {e}")
        raise HTTPException(status_code=500, detail=f"Ошибка получения: {str(e)}")


@app.put("/api/kb/{chunk_id}", response_model=KBUpdateResponse)
async def update_chunk(chunk_id: int, request: KBChunkUpdate):
    """
    Обновление chunk в KB
    """
    try:
        # Формируем updates для metadata
        metadata_updates = {}
        if request.category:
            metadata_updates["category"] = request.category.value
        if request.target_audience:
            metadata_updates["target_audience"] = request.target_audience.value
        if request.priority:
            metadata_updates["priority"] = request.priority.value
        if request.subcategory is not None:
            metadata_updates["subcategory"] = request.subcategory
        if request.tags is not None:
            metadata_updates["tags"] = request.tags
        if request.source is not None:
            metadata_updates["source"] = request.source
        if request.version is not None:
            metadata_updates["version"] = request.version
        if request.related_links is not None:
            metadata_updates["related_links"] = request.related_links
        if request.settlement_id is not None:
            metadata_updates["settlement_id"] = request.settlement_id
        
        if request.use_case or request.stage:
            context = {}
            if request.use_case:
                context["use_case"] = request.use_case
            if request.stage:
                context["stage"] = request.stage
            metadata_updates["context"] = context
        
        if request.metadata:
            metadata_updates.update(request.metadata)
        
        success = await kb_service.update_chunk(
            chunk_id=chunk_id,
            content=request.content,
            metadata_updates=metadata_updates if metadata_updates else None,
            is_active=request.is_active
        )
        
        if not success:
            raise HTTPException(status_code=404, detail=f"Chunk {chunk_id} не найден или не обновлен")
        
        return KBUpdateResponse(
            success=True,
            chunk_id=chunk_id,
            message=f"Chunk {chunk_id} успешно обновлен"
        )
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Ошибка обновления chunk {chunk_id}: {e}")
        raise HTTPException(status_code=500, detail=f"Ошибка обновления: {str(e)}")


@app.delete("/api/kb/{chunk_id}", response_model=KBDeleteResponse)
async def delete_chunk(chunk_id: int, soft_delete: bool = True):
    """
    Удаление chunk из KB
    
    Args:
        chunk_id: ID chunk для удаления
        soft_delete: Если True, помечает is_active=False, иначе физически удаляет
    """
    try:
        success = await kb_service.delete_chunk(chunk_id, soft_delete=soft_delete)
        if not success:
            raise HTTPException(status_code=404, detail=f"Chunk {chunk_id} не найден")
        
        return KBDeleteResponse(
            success=True,
            message=f"Chunk {chunk_id} {'деактивирован' if soft_delete else 'удален'}"
        )
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Ошибка удаления chunk {chunk_id}: {e}")
        raise HTTPException(status_code=500, detail=f"Ошибка удаления: {str(e)}")


@app.post("/api/kb/import", response_model=KBImportResponse)
async def import_kb(request: KBImportRequest):
    """
    Импорт текста в KB с автоматическим разбиением на chunks
    """
    try:
        # Определяем источник текста
        if request.file_path:
            # Читаем из файла
            with open(request.file_path, 'r', encoding='utf-8') as f:
                content = f.read()
        elif request.content:
            content = request.content
        else:
            raise HTTPException(status_code=400, detail="Необходимо указать file_path или content")
        
        # Импортируем (metadata.images / metadata.documents — для релевантного ответа с картинками/документами)
        result = await kb_service.import_from_text(
            text=content,
            category=request.category.value,
            target_audience=request.target_audience.value,
            source=request.source,
            chunk_size=request.chunk_size,
            chunk_overlap=request.chunk_overlap,
            settlement_id=request.settlement_id,
            metadata=request.metadata,
        )
        
        errors = []
        if result["failed"] > 0:
            errors.append(f"Не удалось добавить {result['failed']} chunks")
        
        return KBImportResponse(
            success=result["failed"] == 0,
            chunks_added=result["added"],
            chunks_updated=0,  # При импорте не обновляем существующие
            chunks_failed=result["failed"],
            errors=errors,
            message=f"Импорт завершен: добавлено {result['added']}, ошибок {result['failed']}"
        )
    except HTTPException:
        raise
    except FileNotFoundError:
        raise HTTPException(status_code=404, detail=f"Файл {request.file_path} не найден")
    except Exception as e:
        logger.error(f"Ошибка импорта: {e}")
        import traceback
        logger.error(traceback.format_exc())
        raise HTTPException(status_code=500, detail=f"Ошибка импорта: {str(e)}")


def _extract_media_urls_from_text(text: str) -> Dict[str, Any]:
    """Извлекает из текста URL картинок и документов для metadata.images / metadata.documents."""
    import re
    # URL в тексте (http/https, до пробела/скобки/кавычки/конца строки)
    url_pattern = re.compile(
        r"https?://[^\s\)\]\"\'\>]+",
        re.IGNORECASE,
    )
    seen = set()
    images: List[Dict[str, Any]] = []
    documents: List[Dict[str, Any]] = []
    image_ext = (".jpg", ".jpeg", ".png", ".gif", ".webp", ".bmp")
    doc_ext = (".pdf", ".doc", ".docx", ".xls", ".xlsx", ".pptx")
    for m in url_pattern.finditer(text):
        raw = m.group(0).rstrip(".,;:)")
        if raw in seen:
            continue
        seen.add(raw)
        lower = raw.lower()
        if any(lower.endswith(ext) or ext + "?" in lower or ext + "&" in lower for ext in image_ext):
            images.append({"url": raw, "description": "", "alt": ""})
        elif any(lower.endswith(ext) or ext + "?" in lower or ext + "&" in lower for ext in doc_ext):
            documents.append({"url": raw, "title": "Документ", "description": ""})
    return {"images": images[:20], "documents": documents[:10]}


def _extract_text_from_document(filename: str, content: bytes) -> str:
    """Извлекает текст из файла .txt, .pdf, .docx или .pptx."""
    suf = Path(filename).suffix.lower()
    if suf == ".txt":
        for enc in ("utf-8", "cp1251"):
            try:
                return content.decode(enc)
            except UnicodeDecodeError:
                continue
        raise HTTPException(status_code=400, detail="Не удалось декодировать TXT (попробуйте utf-8 или cp1251)")
    if suf == ".pdf":
        try:
            from pypdf import PdfReader
            import io
            reader = PdfReader(io.BytesIO(content))
            parts = []
            for page in reader.pages:
                t = page.extract_text()
                if t:
                    parts.append(t)
            return "\n\n".join(parts) if parts else ""
        except Exception as e:
            logger.error(f"Ошибка извлечения текста из PDF: {e}")
            raise HTTPException(status_code=400, detail=f"Не удалось прочитать PDF: {e}")
    if suf == ".docx":
        try:
            import io
            from docx import Document
            doc = Document(io.BytesIO(content))
            return "\n\n".join(p.text for p in doc.paragraphs)
        except Exception as e:
            logger.error(f"Ошибка извлечения текста из Word: {e}")
            raise HTTPException(status_code=400, detail=f"Не удалось прочитать Word: {e}")
    if suf == ".pptx":
        try:
            import io
            from pptx import Presentation
        except ImportError as e:
            logger.error("Модуль python-pptx не установлен: %s", e)
            raise HTTPException(
                status_code=400,
                detail="Не удалось прочитать PPTX: установите python-pptx в окружении kb-service: pip install python-pptx"
            )
        try:
            prs = Presentation(io.BytesIO(content))
            parts = []
            for slide in prs.slides:
                slide_texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            if para.text:
                                slide_texts.append(para.text.strip())
                    elif getattr(shape, "text", None):
                        slide_texts.append(shape.text.strip())
                if slide_texts:
                    parts.append("\n".join(slide_texts))
            return "\n\n".join(parts) if parts else ""
        except Exception as e:
            logger.error(f"Ошибка извлечения текста из PowerPoint: {e}")
            raise HTTPException(status_code=400, detail=f"Не удалось прочитать PPTX: {e}")
    raise HTTPException(status_code=400, detail=f"Формат не поддерживается: {suf}. Разрешены .txt, .pdf, .docx, .pptx")


@app.post("/api/kb/import_document", response_model=KBImportResponse)
async def import_document(
    file: UploadFile = File(...),
    category: str = Form("product_info"),
    target_audience: str = Form("both"),
    source: str = Form(""),
    chunk_size: int = Form(3000, ge=100, le=10000),
    chunk_overlap: int = Form(300, ge=0, le=1000),
):
    """
    Импорт документа (TXT, PDF, Word, PowerPoint) в KB: извлечение текста и разбиение на chunks.
    """
    if not file.filename:
        raise HTTPException(status_code=400, detail="Файл не выбран")
    content = await file.read()
    max_mb = int(os.environ.get("KB_IMPORT_MAX_FILE_MB", "100"))
    if len(content) > max_mb * 1024 * 1024:
        raise HTTPException(status_code=400, detail=f"Файл не более {max_mb} МБ")
    try:
        text = _extract_text_from_document(file.filename, content)
    except HTTPException:
        raise
    if not text.strip():
        raise HTTPException(status_code=400, detail="В файле не найден текст")
    try:
        cat_enum = Category(category)
    except ValueError:
        cat_enum = Category.PRODUCT_INFO
    try:
        aud_enum = TargetAudience(target_audience)
    except ValueError:
        aud_enum = TargetAudience.BOTH
    # Извлекаем URL картинок и документов из текста — попадут в metadata для релевантного ответа агента
    extracted = _extract_media_urls_from_text(text)
    import_metadata: Dict[str, Any] = {}
    if extracted.get("images"):
        import_metadata["images"] = extracted["images"]
    if extracted.get("documents"):
        import_metadata["documents"] = extracted["documents"]
    try:
        result = await kb_service.import_from_text(
            text=text,
            category=cat_enum.value,
            target_audience=aud_enum.value,
            source=(source.strip() or file.filename),  # имя файла — источник, если URL не задан
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            metadata=import_metadata if import_metadata else None,
        )
        errors = []
        if result["failed"] > 0:
            errors.append(f"Не удалось добавить {result['failed']} chunks")
        return KBImportResponse(
            success=result["failed"] == 0,
            chunks_added=result["added"],
            chunks_updated=0,
            chunks_failed=result["failed"],
            errors=errors,
            message=f"Импорт документа завершен: добавлено {result['added']}, ошибок {result['failed']}",
        )
    except Exception as e:
        logger.error(f"Ошибка импорта документа: {e}")
        import traceback
        logger.error(traceback.format_exc())
        raise HTTPException(status_code=500, detail=str(e))


# Обработчик ошибок
@app.exception_handler(Exception)
async def global_exception_handler(request, exc):
    """
    Глобальный обработчик ошибок
    """
    logger.error(f"Необработанная ошибка: {exc}")
    import traceback
    logger.error(traceback.format_exc())
    return JSONResponse(
        status_code=500,
        content={
            "error": "Внутренняя ошибка сервера",
            "error_code": "INTERNAL_ERROR",
            "detail": str(exc)
        }
    )


_log_routes()


if __name__ == "__main__":
    import uvicorn
    
    host = os.getenv("API_HOST", "0.0.0.0")
    port = int(os.getenv("API_PORT", "8001"))
    
    uvicorn.run(app, host=host, port=port)

